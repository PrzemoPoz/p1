import csv
from pptx import Presentation


with open("dane/titanic_train.csv") as csvfile:
    data = csv.DictReader(csvfile, delimiter=",")
    sum_m = 0
    sum_f = 0
    i_m = 0
    i_f = 0
    sum_s_m = 0
    sum_s_f = 0
    sum_cl_1 = 0
    sum_cl_2 = 0
    sum_cl_3 = 0
    sum_a_cl_1 = 0
    sum_a_cl_2 = 0
    sum_a_cl_3 = 0
    i_m_np = 0
    i_f_np = 0
    for row in data:
        if row['Sex'] != '':
            if row['Sex'] == 'male':
                if row['Age'] != '':
                    sum_m += float(row['Age'])
                    i_m_np += 1
                i_m += 1
                if row['Survived'] == '1':
                    sum_s_m += 1
            elif row['Sex'] == 'female':
                if row['Age'] != '':
                    sum_f += float(row['Age'])
                    i_f_np += 1
                i_f += 1
                if row['Survived'] == '1':
                    sum_s_f += 1

        if row['Pclass'] == '1':
            sum_a_cl_1 += 1
            if row['Survived'] == '1':
                sum_cl_1 += 1
        elif row['Pclass'] == '2':
            sum_a_cl_2 += 1
            if row['Survived'] == '1':
                sum_cl_2 += 1
        elif row['Pclass'] == '3':
            sum_a_cl_3 += 1
            if row['Survived'] == '1':
                sum_cl_3 += 1

    prs = Presentation()
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Statystyki wg płci i wieku"
    tf = body_shape.text_frame
    tf.text = "Kobiety"
    p = tf.add_paragraph()
    p.text = f'Średni wiek: {round(sum_f / i_f_np, 2)} lat'
    p.level = 1
    p=tf.add_paragraph()
    p.text=f'Przeżywalność: {round(sum_s_f / i_f*100, 2)} %'
    p.level=2

    shapes = slide.shapes
    body_shape1 = shapes.placeholders[2]

    tf1 = body_shape1.text_frame
    tf1.text = "Mężczyźni"
    p = tf1.add_paragraph()
    p.text = f'Średni wiek: {round(sum_m / i_m_np, 2)} lat'
    p.level = 1
    p=tf1.add_paragraph()
    p.text=f'Przeżywalność: {round(sum_s_m / i_m*100, 2)} %'
    p.level=2
    prs.save("raport.pptx")

